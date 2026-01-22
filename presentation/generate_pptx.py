#!/usr/bin/env python3
"""Generate PowerPoint presentation from noodle2 results."""

import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


# Theme colors
BG_COLOR = RGBColor(0x1a, 0x1a, 0x2e)  # Dark blue background
TEXT_COLOR = RGBColor(0xee, 0xee, 0xee)  # Light text
H1_COLOR = RGBColor(0x4a, 0x90, 0xe2)   # Blue for h1
H2_COLOR = RGBColor(0x50, 0xc8, 0x78)   # Green for h2
ACCENT_COLOR = RGBColor(0xff, 0xd7, 0x00)  # Gold for emphasis
DANGER_COLOR = RGBColor(0xe7, 0x4c, 0x3c)  # Red for warnings


def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert hex string to RGBColor."""
    hex_str = hex_str.lstrip('#')
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


def set_slide_background(slide, color: RGBColor):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title: str, subtitle: str = "", author: str = "", date: str = ""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = H1_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = H2_COLOR
        p.alignment = PP_ALIGN.CENTER

    # Author/date
    if author or date:
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
        tf = info_box.text_frame
        if author:
            p = tf.paragraphs[0]
            p.text = author
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_COLOR
            p.alignment = PP_ALIGN.CENTER
        if date:
            p = tf.add_paragraph()
            p.text = date
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_COLOR
            p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title: str, subtitle: str = ""):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = H1_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = H2_COLOR
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title: str, content: list, image_path: str = None, base_path: str = None):
    """Add a content slide with bullet points and/or image."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = H1_COLOR

    content_top = Inches(1.2)
    content_width = Inches(9) if not image_path else Inches(4.5)

    # Content bullets
    if content:
        content_box = slide.shapes.add_textbox(Inches(0.5), content_top, content_width, Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Handle bullet formatting
            text = item.strip()
            if text.startswith('- '):
                text = text[2:]
            elif text.startswith('* '):
                text = text[2:]

            # Check for bold text
            bold_match = re.match(r'\*\*(.+?)\*\*(.*)$', text)
            if bold_match:
                run = p.add_run()
                run.text = bold_match.group(1)
                run.font.bold = True
                run.font.color.rgb = ACCENT_COLOR
                run.font.size = Pt(18)

                if bold_match.group(2):
                    run2 = p.add_run()
                    run2.text = bold_match.group(2)
                    run2.font.color.rgb = TEXT_COLOR
                    run2.font.size = Pt(18)
            else:
                p.text = text
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_COLOR

            p.space_after = Pt(8)

    # Image
    if image_path and base_path:
        full_path = os.path.join(base_path, image_path)
        if os.path.exists(full_path):
            left = Inches(5.2) if content else Inches(1.5)
            top = content_top if content else Inches(1.5)
            width = Inches(4.5) if content else Inches(7)
            try:
                slide.shapes.add_picture(full_path, left, top, width=width)
            except Exception as e:
                print(f"Warning: Could not add image {full_path}: {e}")

    return slide


def add_table_slide(prs, title: str, headers: list, rows: list):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = H1_COLOR

    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1

    table_width = min(Inches(9), Inches(num_cols * 1.8))
    left = (Inches(10) - table_width) / 2

    table = slide.shapes.add_table(num_rows, num_cols, left, Inches(1.3), table_width, Inches(0.4 * num_rows)).table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = H1_COLOR
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2a, 0x2a, 0x40)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_COLOR
            p.alignment = PP_ALIGN.CENTER

    return slide


def add_image_slide(prs, title: str, image_path: str, base_path: str, caption: str = ""):
    """Add a slide with a centered image."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = H1_COLOR

    # Image
    full_path = os.path.join(base_path, image_path)
    if os.path.exists(full_path):
        try:
            # Add centered image
            slide.shapes.add_picture(full_path, Inches(1.25), Inches(1.3), width=Inches(7.5))
        except Exception as e:
            print(f"Warning: Could not add image {full_path}: {e}")

    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER

    return slide


def create_presentation(output_path: str, base_path: str):
    """Create the complete PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Noodle2",
        "Physical Design ECO Orchestration System",
        "Yaakoub Elkhamra, PhD",
        "January 2026"
    )

    # Slide 2: What is Noodle2?
    add_content_slide(prs, "What is Noodle2?", [
        "**ECO Orchestration System** for physical design timing closure",
        "Automatically applies and evaluates **Engineering Change Orders**",
        "Uses **parallel execution** with Ray for trial exploration",
        "Implements **prior learning** to avoid ineffective ECOs",
        "Supports **checkpoint/rollback** for robustness",
        "",
        "All results from REAL OpenROAD execution - no mocking"
    ], base_path=base_path)

    # Slide 3: Technology Stack
    add_table_slide(prs, "Technology Stack",
        ["Layer", "Technology", "Purpose"],
        [
            ["Orchestration", "Noodle2 (Python)", "ECO selection, survivor management"],
            ["Parallelism", "Ray", "Distributed trial execution"],
            ["EDA Engine", "OpenROAD", "Physical design, timing analysis"],
            ["Build Flow", "ORFS", "Synthesis, placement, routing"],
            ["PDKs", "Nangate45, ASAP7, Sky130", "Process design kits"],
        ]
    )

    # Slide 4: OpenROAD
    add_content_slide(prs, "OpenROAD - The EDA Engine", [
        "**OpenSTA** - Static Timing Analysis (WNS, TNS, hot_ratio)",
        "**TritonPlace** - Global and detailed placement",
        "**FastRoute** - Global routing",
        "**TritonRoute** - Detailed routing",
        "**ReSizer** - Gate sizing, buffer insertion",
        "",
        "Production-quality, reproducible, open-source"
    ], base_path=base_path)

    # Slide 5: Ray
    add_content_slide(prs, "Ray - Distributed Parallel Execution", [
        "**Single-Node Mode:** 25 trials per stage in parallel",
        "Utilizes all CPU cores (32 cores in demo)",
        "Shared memory for ODB file access",
        "",
        "**Multi-Node Mode:** Linear scaling with cluster",
        "Fault tolerance for failed trials",
        "Dashboard monitoring at http://localhost:8265"
    ], base_path=base_path)

    # Slide 6: ECO Types
    add_table_slide(prs, "ECO Types Supported",
        ["Category", "ECOs"],
        [
            ["Topology Neutral", "Cell Resize, Buffer Insert/Remove, Pin Swap"],
            ["Placement Affecting", "Timing-Driven Placement, Density Opt"],
            ["Global/Aggressive", "Full Optimization, Multi-Pass, VT Swap"],
            ["Repair", "Hold Repair, Clock Net Repair, Tie Fanout"],
        ]
    )

    # Slide 7: Study Flow
    add_content_slide(prs, "Study Execution Flow", [
        "1. **Base Case Verification** - Validate initial design metrics",
        "2. **Stage Execution** - 20 stages with 25 trials each",
        "3. **Survivor Selection** - Keep best performing variants",
        "4. **Prior Learning** - Track ECO effectiveness",
        "5. **Checkpoint/Rollback** - Save state, recover from degradation",
        "6. **Visualization** - Generate heatmaps, trajectories",
        "",
        "**Total: 500 trials per study**"
    ], base_path=base_path)

    # Slide 8: Nangate45 Section
    add_section_slide(prs, "Nangate45", "45nm Educational PDK")

    # Slide 9: Nangate45 Setup
    add_table_slide(prs, "Nangate45 - Design Setup",
        ["Property", "Value"],
        [
            ["PDK", "Nangate45 (45nm)"],
            ["Design", "Ibex RISC-V Core"],
            ["Cell Count", "~10,000 cells"],
            ["Clock Period", "Aggressive constraints"],
            ["ODB Size", "8.3 MB"],
        ]
    )

    # Slide 10: Nangate45 Stage Progression
    add_image_slide(prs, "Nangate45 - Stage Progression",
        "images/nangate45/stage_progression.png", base_path,
        "20 stages, funnel survivor selection (8 -> 2)")

    # Slide 11: Nangate45 WNS
    add_image_slide(prs, "Nangate45 - WNS Trajectory",
        "images/nangate45/wns_trajectory.png", base_path)

    # Slide 12: Nangate45 Hot Ratio
    add_image_slide(prs, "Nangate45 - Hot Ratio Trajectory",
        "images/nangate45/hot_ratio_trajectory.png", base_path)

    # Slide 13: Nangate45 Results
    add_table_slide(prs, "Nangate45 - Final Results",
        ["Metric", "Initial", "Final", "Improvement"],
        [
            ["WNS", "-1848 ps", "-1576 ps", "+14.7%"],
            ["hot_ratio", "0.523", "0.149", "-71.6%"],
        ]
    )

    # Slide 14: ASAP7 Section
    add_section_slide(prs, "ASAP7", "7nm Predictive PDK")

    # Slide 15: ASAP7 Setup
    add_table_slide(prs, "ASAP7 - Design Setup",
        ["Property", "Value"],
        [
            ["PDK", "ASAP7 (7nm predictive)"],
            ["Design", "Ibex RISC-V Core"],
            ["Cell Count", "~10,000 cells"],
            ["Clock Period", "Aggressive 7nm timing"],
            ["ODB Size", "~10 MB"],
        ]
    )

    # Slide 16: ASAP7 Stage Progression
    add_image_slide(prs, "ASAP7 - Stage Progression",
        "images/asap7/stage_progression.png", base_path)

    # Slide 17: ASAP7 WNS
    add_image_slide(prs, "ASAP7 - WNS Trajectory",
        "images/asap7/wns_trajectory.png", base_path)

    # Slide 18: ASAP7 Results
    add_table_slide(prs, "ASAP7 - Final Results",
        ["Metric", "Initial", "Final", "Improvement"],
        [
            ["WNS", "-1067 ps", "-1004 ps", "+5.9%"],
            ["hot_ratio", "0.55", "0.004", "-99.3%"],
        ]
    )

    # Slide 19: Sky130 Section
    add_section_slide(prs, "Sky130 + Microwatt", "130nm Open PDK + OpenPOWER Core")

    # Slide 20: Sky130 Setup
    add_table_slide(prs, "Sky130 Microwatt - Design Setup",
        ["Property", "Value"],
        [
            ["PDK", "Sky130HD (130nm)"],
            ["Design", "Microwatt OpenPOWER Core"],
            ["Cell Count", "162,637 cells"],
            ["Clock Period", "4ns (250 MHz)"],
            ["ODB Size", "95 MB"],
        ]
    )

    # Slide 21: Sky130 Stage Progression
    add_image_slide(prs, "Sky130 Microwatt - Stage Progression",
        "images/sky130/stage_progression.png", base_path,
        "Stage 6 (red): degradation detected but below rollback threshold")

    # Slide 22: Sky130 WNS
    add_image_slide(prs, "Sky130 Microwatt - WNS Trajectory",
        "images/sky130/wns_trajectory.png", base_path)

    # Slide 23: Sky130 Hot Ratio
    add_image_slide(prs, "Sky130 Microwatt - Hot Ratio Trajectory",
        "images/sky130/hot_ratio_trajectory.png", base_path,
        "hot_ratio: 0.097 -> 0.0003 (99.7% reduction)")

    # Slide 24: Sky130 Results
    add_table_slide(prs, "Sky130 Microwatt - Final Results",
        ["Metric", "Initial", "Final", "Improvement"],
        [
            ["WNS", "-2989 ps", "-1466 ps", "+51.0%"],
            ["hot_ratio", "0.097", "0.0003", "-99.7%"],
        ]
    )

    # Slide 25: Long Pole Analysis
    add_content_slide(prs, "The 'Long Pole in the Tent' Problem", [
        "**Observation:** hot_ratio improved 99.7% but WNS only 51%",
        "",
        "**hot_ratio:** Fraction of paths violating timing (9.7% -> 0.03%)",
        "**WNS:** Worst single path's negative slack (-2989ps -> -1466ps)",
        "",
        "**Meaning:**",
        "- We fixed 99.7% of all timing violations",
        "- But the single worst path only improved 51%",
        "- ONE stubborn path dominates WNS"
    ], base_path=base_path)

    # Slide 26: Why Worst Path Resists
    add_content_slide(prs, "Why Does the Worst Path Resist?", [
        "**Long combinational depth** - Many logic stages in series",
        "**Placement-limited** - Cells far apart, wire delay dominates",
        "**Technology limits** - Already using max drive strength",
        "**Critical logic** - Carry chains, multipliers, barrel shifters",
        "**Routing congestion** - Detours adding wire delay",
        "",
        "ECOs fix 'moderately bad' paths but struggle with 'truly terrible' ones"
    ], base_path=base_path)

    # Slide 27: Cross-PDK Comparison
    add_table_slide(prs, "Cross-PDK Comparison",
        ["PDK", "Design", "Cells", "WNS Impr.", "hot_ratio Red."],
        [
            ["Nangate45", "Ibex", "~10K", "+14.7%", "-71.6%"],
            ["ASAP7", "Ibex", "~10K", "+5.9%", "-99.3%"],
            ["Sky130", "Microwatt", "162K", "+51.0%", "-99.7%"],
        ]
    )

    # Slide 28: ECO Leaderboard
    add_table_slide(prs, "ECO Effectiveness Leaderboard",
        ["ECO", "Applications", "Success", "Rate"],
        [
            ["hold_repair", "101", "101", "100%"],
            ["aggressive_timing", "99", "99", "100%"],
            ["multi_pass_timing", "97", "97", "100%"],
            ["buffer_insertion", "77", "77", "100%"],
            ["gate_cloning", "76", "0", "0%"],
            ["dead_logic_elim", "72", "0", "0%"],
        ]
    )

    # Slide 29: ECO Categories
    add_table_slide(prs, "ECO Success by Category",
        ["Category", "ECOs", "Success"],
        [
            ["Repair", "hold, sequential, design, tie_fanout", "100%"],
            ["Timing", "aggressive, multi_pass, full_opt", "100%"],
            ["Cell", "resize, swap, buffer insert/remove", "100%"],
            ["Placement", "timing_driven, density, iterative", "100%"],
            ["Failed", "gate_cloning, dead_logic_elim", "0%"],
        ]
    )

    # Slide 30: Key Features
    add_content_slide(prs, "Key Noodle2 Features Demonstrated", [
        "**Parallel Execution** - Ray-based trial parallelism",
        "**Prior Learning** - ECO effectiveness tracking",
        "**Checkpoint/Rollback** - Recovery from degradation",
        "**Degradation Detection** - Stage 6 flagged without rollback",
        "**Multi-PDK Support** - Nangate45, ASAP7, Sky130",
        "**Safety Domains** - Sandbox, guarded, locked modes"
    ], base_path=base_path)

    # Slide 31: Conclusions
    add_content_slide(prs, "Conclusions", [
        "Noodle2 improved timing on **3 different PDKs**",
        "**162K cell Microwatt** achieved 51% WNS, 99.7% violation reduction",
        "All executions used **real OpenROAD** - no simulation",
        "Prior learning effectively filters ineffective ECOs",
        "System scales from 10K to 160K+ cells",
        "",
        "**Key insight:** ECOs excel at fixing most violations;",
        "worst path(s) may require architectural changes",
        "",
        "**Total: 1,500 real ECO trials executed**"
    ], base_path=base_path)

    # Slide 32: Thank You
    add_title_slide(
        prs,
        "Thank You",
        "Noodle2 - Physical Design ECO Orchestration",
        "Yaakoub Elkhamra, PhD",
        "All visualizations from actual OpenROAD execution data"
    )

    # Save
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    base_path = Path(__file__).parent / "bundle"
    output_path = Path(__file__).parent / "bundle" / "noodle2_results.pptx"
    create_presentation(str(output_path), str(base_path))
